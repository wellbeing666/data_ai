import json
from pathlib import Path
from typing import Any

from fastapi import HTTPException, status


REPORT_FILENAME = "report.md"
PPTX_FILENAME = "report.pptx"
PPTX_PREVIEW_FILENAME = "pptx_preview.json"


def generate_markdown_report(
    analysis_result_path: str,
    chart_paths: list[str] | None = None,
    output_path: str | None = None,
    include_pptx: bool = True,
) -> dict[str, Any]:
    result_path = Path(analysis_result_path).resolve()
    if not result_path.exists() or not result_path.is_file():
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="分析结果文件不存在，请先完成分析后再生成报告。",
        )

    analysis_result = _load_json(result_path)
    result_dir = result_path.parent
    resolved_output_path = Path(output_path).resolve() if output_path else result_dir / REPORT_FILENAME

    if not _is_relative_to(resolved_output_path, result_dir):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="报告文件必须保存到当前任务目录内。",
        )

    resolved_chart_paths = _resolve_chart_paths(
        chart_paths=chart_paths or [],
        analysis_result=analysis_result,
        base_dir=result_dir,
    )
    explanation = _load_optional_json(result_dir / "explanation.json") or _load_optional_json(result_dir / "prediction_explanation.json")
    if isinstance(explanation, dict) and _has_explanation_content(explanation):
        markdown = _build_explanation_report(explanation, analysis_result, resolved_chart_paths, result_dir)
        markdown = _append_cleaning_section(markdown, result_dir)
        resolved_output_path.write_text(markdown, encoding="utf-8")
        return _build_report_response(
            report_path=resolved_output_path,
            analysis_result_path=result_path,
            chart_paths=resolved_chart_paths,
            explanation=explanation,
            analysis_result=analysis_result,
            include_pptx=include_pptx,
        )

    task_type = str(
        analysis_result.get("task_type")
        or analysis_result.get("analysis_type")
        or _deep_get(analysis_result, ["analysis_plan", "task_type"])
        or ""
    )
    if task_type == "grade_analysis":
        markdown = _build_grade_analysis_report(analysis_result, resolved_chart_paths, result_dir)
    elif task_type == "sales_decline_analysis":
        markdown = _build_sales_decline_report(analysis_result, resolved_chart_paths, result_dir)
    else:
        markdown = _build_general_report(analysis_result, resolved_chart_paths, result_dir)

    markdown = _append_cleaning_section(markdown, result_dir)
    resolved_output_path.write_text(markdown, encoding="utf-8")
    return _build_report_response(
        report_path=resolved_output_path,
        analysis_result_path=result_path,
        chart_paths=resolved_chart_paths,
        explanation=explanation if isinstance(explanation, dict) else {},
        analysis_result=analysis_result,
        include_pptx=include_pptx,
    )


def generate_pptx_report(
    analysis_result_path: str,
    chart_paths: list[str] | None = None,
) -> dict[str, Any]:
    result_path = Path(analysis_result_path).resolve()
    if not result_path.exists() or not result_path.is_file():
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="分析结果文件不存在，请先完成分析后再生成 PPTX。",
        )

    analysis_result = _load_json(result_path)
    result_dir = result_path.parent
    resolved_chart_paths = _resolve_chart_paths(
        chart_paths=chart_paths or [],
        analysis_result=analysis_result,
        base_dir=result_dir,
    )
    explanation = _load_optional_json(result_dir / "explanation.json") or _load_optional_json(result_dir / "prediction_explanation.json")
    if not isinstance(explanation, dict):
        explanation = {}

    pptx_path = _create_pptx_report(result_dir, explanation, analysis_result, resolved_chart_paths)
    if pptx_path is None:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="PPTX 生成失败，请稍后重试。",
        )
    preview_path = _create_pptx_preview(result_dir, explanation, analysis_result, resolved_chart_paths, pptx_path)
    return {
        "analysis_result_path": str(result_path),
        "chart_paths": [str(path) for path in resolved_chart_paths],
        "pptx_path": str(pptx_path),
        "pptx_preview_path": str(preview_path),
    }

def _load_json(path: Path) -> dict[str, Any]:
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"analysis_result.json is not valid JSON: {exc}",
        ) from exc

    if not isinstance(data, dict):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="analysis_result.json must be a JSON object.",
        )

    return data


def _load_optional_json(path: Path) -> Any | None:
    if not path.exists() or not path.is_file():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return None


def _resolve_chart_paths(
    chart_paths: list[str],
    analysis_result: dict[str, Any],
    base_dir: Path,
) -> list[Path]:
    raw_paths = chart_paths or [str(path) for path in analysis_result.get("charts", [])]
    resolved_paths: list[Path] = []
    seen: set[str] = set()

    for raw_path in raw_paths:
        path = _resolve_chart_path(raw_path, base_dir)
        if path is None:
            continue
        key = path.as_posix()
        if key not in seen:
            resolved_paths.append(path)
            seen.add(key)

    charts_dir = base_dir / "charts"
    if charts_dir.exists():
        for chart_file in sorted(charts_dir.glob("*.png")):
            path = chart_file.resolve()
            key = path.as_posix()
            if path.is_file() and key not in seen:
                resolved_paths.append(path)
                seen.add(key)

    return resolved_paths


def _resolve_chart_path(raw_path: Any, base_dir: Path) -> Path | None:
    raw = str(raw_path or "").strip().replace("\\", "/")
    if not raw or raw.startswith(("http://", "https://")):
        return None

    if raw.startswith("/storage/"):
        candidate = Path(raw.lstrip("/"))
    elif raw.startswith("storage/"):
        candidate = Path(raw)
    else:
        path = Path(raw)
        if path.is_absolute():
            candidate = path
        elif raw.startswith("charts/"):
            candidate = base_dir / raw
        elif "/" not in raw:
            candidate = base_dir / "charts" / raw
        else:
            candidate = base_dir / raw

    candidate = candidate.resolve()
    if candidate.exists() and candidate.is_file():
        return candidate
    return None

def _build_grade_analysis_report(
    analysis_result: dict[str, Any],
    chart_paths: list[Path],
    report_dir: Path,
) -> str:
    summary = _as_list(analysis_result.get("summary"))
    task_name = "成绩按班级统计分析报告"
    thresholds = analysis_result.get("thresholds", {})
    pass_score = thresholds.get("pass_score", 60)
    excellent_score = thresholds.get("excellent_score", 90)

    total_students = sum(_safe_number(row.get("student_count")) for row in summary)
    class_count = len(summary)
    best_average = _max_row(summary, "average_score")
    lowest_average = _min_row(summary, "average_score")
    best_pass_rate = _max_row(summary, "pass_rate")
    best_excellent_rate = _max_row(summary, "excellent_rate")

    lines = [
        f"# {task_name}",
        "",
        "## 分析目标",
        "按班级统计学生成绩，展示班级之间的平均分、及格率和优秀率差异，为课程设计演示提供结构化分析结果。",
        "",
        "## 数据概况",
        f"- 班级数量：{class_count}",
        f"- 有效成绩记录数：{int(total_students)}",
        f"- 及格线：{pass_score} 分",
        f"- 优秀线：{excellent_score} 分",
        "",
        "## 关键发现",
    ]

    if best_average:
        lines.append(
            f"- 平均分最高的班级是 **{best_average.get('class_name')}**，平均分为 **{_fmt_number(best_average.get('average_score'))}**。"
        )
    if lowest_average:
        lines.append(
            f"- 平均分最低的班级是 **{lowest_average.get('class_name')}**，平均分为 **{_fmt_number(lowest_average.get('average_score'))}**。"
        )
    if best_pass_rate:
        lines.append(
            f"- 及格率最高的班级是 **{best_pass_rate.get('class_name')}**，及格率为 **{_fmt_percent(best_pass_rate.get('pass_rate'))}**。"
        )
    if best_excellent_rate:
        lines.append(
            f"- 优秀率最高的班级是 **{best_excellent_rate.get('class_name')}**，优秀率为 **{_fmt_percent(best_excellent_rate.get('excellent_rate'))}**。"
        )
    if not summary:
        lines.append("- 暂未发现可展示的班级统计结果。")

    lines.extend(
        [
            "",
            "## 班级对比明细",
            "",
            "| 班级 | 人数 | 平均分 | 最高分 | 最低分 | 及格率 | 优秀率 |",
            "| --- | ---: | ---: | ---: | ---: | ---: | ---: |",
        ]
    )
    for row in summary:
        lines.append(
            "| {class_name} | {student_count} | {average_score} | {max_score} | {min_score} | {pass_rate} | {excellent_rate} |".format(
                class_name=row.get("class_name", ""),
                student_count=int(_safe_number(row.get("student_count"))),
                average_score=_fmt_number(row.get("average_score")),
                max_score=_fmt_number(row.get("max_score")),
                min_score=_fmt_number(row.get("min_score")),
                pass_rate=_fmt_percent(row.get("pass_rate")),
                excellent_rate=_fmt_percent(row.get("excellent_rate")),
            )
        )

    lines.extend(
        [
            "",
            "## 图表说明",
            *_build_chart_lines(chart_paths, report_dir),
            "",
            "## 可能原因 / 统计结论",
            "- 班级平均分可以反映各班整体成绩水平，适合用于横向对比。",
            "- 及格率体现基础达标情况，优秀率体现高分段学生占比；两者结合能更完整地观察班级成绩结构。",
            "- 如果某个班级平均分不低但优秀率偏低，可能说明学生集中在中等分段；如果及格率偏低，则需要重点关注临界和低分学生。",
            "",
            "## 数据限制",
            "- 当前报告基于上传文件中的一次成绩数据，不能代表长期趋势。",
            "- 及格率和优秀率默认按固定阈值计算，如学校规则不同需要调整。",
            "- 缺失或无法转换为数值的成绩记录会影响统计结果，建议在正式分析前确认数据完整性。",
        ]
    )

    return "\n".join(lines) + "\n"


def _build_sales_decline_report(
    analysis_result: dict[str, Any],
    chart_paths: list[Path],
    report_dir: Path,
) -> str:
    task_name = _deep_get(analysis_result, ["analysis_plan", "task_name"]) or "销量下降分析报告"
    summary = analysis_result.get("summary") or analysis_result.get("result") or []

    lines = [
        f"# {task_name}",
        "",
        "## 分析目标",
        "围绕销量或销售指标的下降现象进行描述性拆解，识别与下降同步变化的时间、渠道、产品或地区因素。",
        "",
        "## 数据概况",
        f"- 分析结果记录数：{len(summary) if isinstance(summary, list) else 1}",
        "- 本报告仅基于已生成的结构化分析结果和图表进行总结。",
        "",
        "## 关键发现",
        *_build_generic_findings(analysis_result),
        "",
        "## 图表说明",
        *_build_chart_lines(chart_paths, report_dir),
        "",
        "## 可能原因 / 统计结论",
        "- 以下内容仅表示统计相关或同步变化，不构成绝对因果判断。",
        "- 如果某些渠道、地区或产品在销量下降期间变化更明显，它们可能是需要优先排查的方向。",
        "- 若时间趋势中存在集中下滑区间，建议结合促销、库存、价格、流量和外部活动进一步验证。",
        "- 当前分析不能单独证明“某因素导致销量下降”，只能为后续业务排查提供线索。",
        "",
        "## 数据限制",
        "- 如果缺少历史同期数据，无法充分判断同比变化。",
        "- 如果缺少渠道、地区、产品或流量字段，下降拆解会受到限制。",
        "- 统计相关不等于因果关系，重要结论需要结合业务背景和实验验证。",
    ]
    return "\n".join(lines) + "\n"


def _build_general_report(
    analysis_result: dict[str, Any],
    chart_paths: list[Path],
    report_dir: Path,
) -> str:
    task_name = _deep_get(analysis_result, ["analysis_plan", "task_name"]) or "数据分析报告"
    lines = [
        f"# {task_name}",
        "",
        "## 分析目标",
        "根据已生成的分析结果，对数据进行结构化总结和展示。",
        "",
        "## 数据概况",
        f"- 任务类型：{analysis_result.get('task_type', 'unknown')}",
        f"- 分析是否成功：{analysis_result.get('success', 'unknown')}",
        "",
        "## 关键发现",
        *_build_generic_findings(analysis_result),
        "",
        "## 图表说明",
        *_build_chart_lines(chart_paths, report_dir),
        "",
        "## 可能原因 / 统计结论",
        "- 当前报告基于已有统计结果进行归纳，适合用于展示主要数据模式。",
        "- 若需要进一步解释变化原因，建议补充时间、维度和业务背景字段。",
        "",
        "## 数据限制",
        "- 当前报告依赖 analysis_result.json 中已有的结果字段。",
        "- 如果原始数据存在缺失、重复或字段含义不清，结论可靠性会受到影响。",
    ]
    return "\n".join(lines) + "\n"


def _build_explanation_report(
    explanation: dict[str, Any],
    analysis_result: dict[str, Any],
    chart_paths: list[Path],
    report_dir: Path,
) -> str:
    title = (
        _deep_get(analysis_result, ["analysis_plan", "analysis_goal"])
        or _deep_get(analysis_result, ["analysis_plan", "task_name"])
        or explanation.get("title")
        or "数据分析报告"
    )
    lines = [
        f"# {title}",
        "",
        "## 摘要",
        str(explanation.get("summary") or "暂无摘要。"),
        "",
        "## 关键发现",
        *_build_explanation_items(explanation.get("key_findings")),
        "",
        "## 建议动作",
        *_build_explanation_items(explanation.get("recommendations")),
        "",
        "## 图表说明",
        *_build_chart_lines(chart_paths, report_dir),
        "",
        "## 限制说明",
        *_build_explanation_items(explanation.get("limitations")),
    ]
    outline = _as_outline(explanation.get("ppt_outline"))
    if outline:
        lines.extend(["", "## PPT 大纲"])
        for index, slide in enumerate(outline, start=1):
            lines.append(f"### {index}. {slide.get('title', '未命名页面')}")
            for bullet in _as_string_list(slide.get("bullets")):
                lines.append(f"- {bullet}")
            chart = str(slide.get("chart") or "")
            if chart:
                lines.append(f"- 关联图表：{Path(chart).name}")
    return "\n".join(lines) + "\n"


def _has_explanation_content(explanation: dict[str, Any]) -> bool:
    return any(
        explanation.get(key)
        for key in ("summary", "key_findings", "recommendations", "limitations", "ppt_outline")
    )


def _build_explanation_items(value: Any) -> list[str]:
    items = _as_string_list(value)
    if not items:
        return ["- 暂无。"]
    return [f"- {item}" for item in items]


def _as_string_list(value: Any) -> list[str]:
    if not isinstance(value, list):
        return []
    result = []
    for item in value:
        if isinstance(item, dict):
            title = item.get("title") or item.get("name") or item.get("finding")
            description = item.get("description") or item.get("text") or item.get("insight")
            if title and description:
                result.append(f"{title}：{description}")
            elif title:
                result.append(str(title))
            elif description:
                result.append(str(description))
        elif item is not None:
            result.append(str(item))
    return result


def _as_outline(value: Any) -> list[dict[str, Any]]:
    if not isinstance(value, list):
        return []
    return [item for item in value if isinstance(item, dict)]


def _build_chart_lines(chart_paths: list[Path], report_dir: Path) -> list[str]:
    if not chart_paths:
        return ["- 暂无可用图表。"]

    lines = []
    for index, chart_path in enumerate(chart_paths, start=1):
        title = _guess_chart_title(chart_path)
        display_path = _markdown_path(chart_path, report_dir)
        lines.extend(
            [
                f"### 图表 {index}：{title}",
                f"![{title}]({display_path})",
                f"- 说明：该图用于展示{title}。",
                "",
            ]
        )
    return lines


def _build_generic_findings(analysis_result: dict[str, Any]) -> list[str]:
    findings = analysis_result.get("key_findings") or analysis_result.get("findings")
    if isinstance(findings, list) and findings:
        lines = []
        for finding in findings[:5]:
            if isinstance(finding, dict):
                title = finding.get("title") or finding.get("name") or "发现"
                description = finding.get("description") or finding.get("insight") or ""
                lines.append(f"- **{title}**：{description}")
            else:
                lines.append(f"- {finding}")
        return lines

    summary = analysis_result.get("summary")
    if isinstance(summary, list) and summary:
        return [f"- 已生成 {len(summary)} 条结构化统计结果，可查看 analysis_result.json 获取明细。"]
    if isinstance(summary, str) and summary:
        return [f"- {summary}"]
    return ["- 暂未发现可自动提取的关键发现。"]


def _guess_chart_title(chart_path: Path) -> str:
    name = chart_path.stem.lower()
    if "average" in name or "avg" in name:
        return "班级平均分"
    if "pass_rate" in name or "pass" in name:
        return "班级及格率"
    if "excellent" in name:
        return "班级优秀率"
    if "trend" in name:
        return "趋势变化"
    return chart_path.stem.replace("_", " ")


def _markdown_path(chart_path: Path, report_dir: Path) -> str:
    try:
        return chart_path.relative_to(report_dir).as_posix()
    except ValueError:
        return chart_path.as_posix()


def _as_list(value: Any) -> list[dict[str, Any]]:
    if isinstance(value, list):
        return [item for item in value if isinstance(item, dict)]
    return []


def _max_row(rows: list[dict[str, Any]], key: str) -> dict[str, Any] | None:
    valid_rows = [row for row in rows if _is_number(row.get(key))]
    return max(valid_rows, key=lambda row: float(row[key]), default=None)


def _min_row(rows: list[dict[str, Any]], key: str) -> dict[str, Any] | None:
    valid_rows = [row for row in rows if _is_number(row.get(key))]
    return min(valid_rows, key=lambda row: float(row[key]), default=None)


def _safe_number(value: Any) -> float:
    if _is_number(value):
        return float(value)
    return 0.0


def _is_number(value: Any) -> bool:
    return isinstance(value, (int, float)) and not isinstance(value, bool)


def _fmt_number(value: Any) -> str:
    if not _is_number(value):
        return "-"
    number = float(value)
    return str(int(number)) if number.is_integer() else f"{number:.2f}"


def _fmt_percent(value: Any) -> str:
    if not _is_number(value):
        return "-"
    return f"{float(value) * 100:.2f}%"


def _deep_get(data: dict[str, Any], keys: list[str]) -> Any:
    value: Any = data
    for key in keys:
        if not isinstance(value, dict):
            return None
        value = value.get(key)
    return value


def _is_relative_to(path: Path, parent: Path) -> bool:
    try:
        path.relative_to(parent)
        return True
    except ValueError:
        return False



def _append_cleaning_section(markdown: str, result_dir: Path) -> str:
    cleaning_report = _load_optional_json(result_dir / "cleaning_report.json")
    if not isinstance(cleaning_report, dict):
        return markdown
    strategies = cleaning_report.get("applied_strategies")
    lines = [markdown.rstrip(), "", "## 数据清洗说明"]
    message = str(cleaning_report.get("message") or "已完成数据质量检查。")
    lines.append(f"- {message}")
    lines.append(f"- 清洗前记录数：{cleaning_report.get('row_count_before', '-')}")
    lines.append(f"- 清洗后记录数：{cleaning_report.get('row_count_after', '-')}")
    if isinstance(strategies, list) and strategies:
        for item in strategies:
            if not isinstance(item, dict):
                continue
            column = item.get("column") or "全表"
            label = item.get("strategy_label") or item.get("strategy_id") or "已确认"
            description = item.get("description") or ""
            lines.append(f"- {column}：{label}。{description}")
    else:
        lines.append("- 未应用自动修改策略，分析基于当前数据版本执行。")
    return "\n".join(lines) + "\n"


def _build_report_response(
    *,
    report_path: Path,
    analysis_result_path: Path,
    chart_paths: list[Path],
    explanation: dict[str, Any],
    analysis_result: dict[str, Any],
    include_pptx: bool = True,
) -> dict[str, Any]:
    pptx_path = None
    pptx_preview_path = None
    if include_pptx:
        pptx_path = _create_pptx_report(report_path.parent, explanation, analysis_result, chart_paths)
        if pptx_path is not None:
            pptx_preview_path = _create_pptx_preview(report_path.parent, explanation, analysis_result, chart_paths, pptx_path)
    return {
        "report_path": str(report_path),
        "analysis_result_path": str(analysis_result_path),
        "chart_paths": [str(path) for path in chart_paths],
        "pptx_path": str(pptx_path) if pptx_path else None,
        "pptx_preview_path": str(pptx_preview_path) if pptx_preview_path else None,
    }


def _create_pptx_report(
    result_dir: Path,
    explanation: dict[str, Any],
    analysis_result: dict[str, Any],
    chart_paths: list[Path],
) -> Path | None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except Exception:
        return None

    slides = _ppt_slide_payload(explanation, analysis_result, chart_paths)
    pptx_path = result_dir / PPTX_FILENAME
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    total = max(len(slides), 1)

    def rgb(hex_value: str) -> RGBColor:
        value = hex_value.strip().lstrip("#")
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

    def add_shape(slide: Any, shape_type: Any, left: float, top: float, width: float, height: float, fill: str, line: str | None = None) -> Any:
        shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if line:
            shape.line.color.rgb = rgb(line)
        else:
            try:
                shape.line.fill.background()
            except Exception:
                pass
        return shape

    def add_text(
        slide: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        *,
        font_size: int,
        color: str = "0F172A",
        bold: bool = False,
        align: Any | None = None,
    ) -> Any:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.02)
        frame.margin_right = Inches(0.02)
        paragraph = frame.paragraphs[0]
        if align is not None:
            paragraph.alignment = align
        run = paragraph.add_run()
        run.text = str(text)
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return box

    def add_bullets(slide: Any, bullets: list[str], left: float, top: float, width: float, height: float, font_size: int = 15) -> Any:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.12)
        frame.margin_right = Inches(0.12)
        frame.margin_top = Inches(0.08)
        frame.margin_bottom = Inches(0.08)
        values = bullets or ["暂无要点。"]
        for idx, bullet in enumerate(values[:7]):
            paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0
            paragraph.space_after = Pt(8)
            for run in paragraph.runs:
                run.font.name = "Microsoft YaHei"
                run.font.size = Pt(font_size)
                run.font.color.rgb = rgb("1F2937")
        return box

    def add_picture(slide: Any, chart_path: Path, left: float, top: float, width: float, height: float) -> None:
        try:
            picture = slide.shapes.add_picture(str(chart_path), Inches(left), Inches(top), width=Inches(width))
            if picture.height > Inches(height):
                scale = Inches(height) / picture.height
                picture.width = int(picture.width * scale)
                picture.height = int(picture.height * scale)
            picture.left = int(Inches(left) + (Inches(width) - picture.width) / 2)
            picture.top = int(Inches(top) + (Inches(height) - picture.height) / 2)
        except Exception:
            add_text(slide, left + 0.2, top + 0.2, width - 0.4, 0.6, "图表暂不可嵌入", font_size=13, color="64748B")

    for index, slide_data in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, "F8FAFC")
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.34, "0F766E")
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 7.22, 13.333, 0.28, "E2E8F0")
        add_text(slide, 0.62, 7.24, 9.0, 0.18, "AI 原生数据分析工作台", font_size=8, color="475569")
        add_text(slide, 11.8, 7.24, 0.9, 0.18, f"{index}/{total}", font_size=8, color="475569", align=PP_ALIGN.RIGHT)

        title = str(slide_data.get("title") or f"第 {index} 页")
        bullets = _as_string_list(slide_data.get("bullets"))[:7]
        subtitle = str(slide_data.get("subtitle") or slide_data.get("section_label") or "数据分析报告")
        chart_path = _slide_chart_path(slide_data, chart_paths)

        if index == 1:
            add_shape(slide, MSO_SHAPE.RECTANGLE, 0.65, 0.95, 5.25, 0.16, "14B8A6")
            add_text(slide, 0.65, 1.25, 8.2, 0.75, title, font_size=32, color="0F172A", bold=True)
            add_text(slide, 0.68, 2.1, 5.7, 0.35, subtitle, font_size=13, color="0F766E", bold=True)
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.68, 2.82, 5.7, 2.9, "FFFFFF", "DDE7EF")
            add_bullets(slide, bullets[:4], 0.95, 3.05, 5.15, 2.35, 16)
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.0, 1.0, 5.55, 4.9, "ECFDF5", "99F6E4")
            if chart_path and chart_path.exists():
                add_picture(slide, chart_path, 7.25, 1.28, 5.05, 4.35)
            else:
                add_text(slide, 7.5, 2.6, 4.6, 0.7, "分析结果已整理为结构化报告", font_size=20, color="0F766E", bold=True, align=PP_ALIGN.CENTER)
        else:
            add_shape(slide, MSO_SHAPE.OVAL, 0.62, 0.82, 0.52, 0.52, "0F766E")
            add_text(slide, 0.755, 0.94, 0.25, 0.18, str(index), font_size=12, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, 1.25, 0.76, 10.8, 0.55, title, font_size=25, color="0F172A", bold=True)
            add_text(slide, 1.28, 1.28, 6.8, 0.28, subtitle, font_size=10, color="64748B")
            if chart_path and chart_path.exists():
                add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 1.82, 5.65, 4.75, "FFFFFF", "DDE7EF")
                add_bullets(slide, bullets, 0.95, 2.05, 5.15, 4.2, 15)
                add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.75, 1.82, 5.85, 4.75, "FFFFFF", "DDE7EF")
                add_picture(slide, chart_path, 6.95, 2.05, 5.45, 4.25)
            else:
                add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 1.82, 11.9, 4.75, "FFFFFF", "DDE7EF")
                add_bullets(slide, bullets, 1.05, 2.1, 11.2, 4.15, 17)

    try:
        prs.save(pptx_path)
    except Exception:
        return None
    return pptx_path


def _create_pptx_preview(
    result_dir: Path,
    explanation: dict[str, Any],
    analysis_result: dict[str, Any],
    chart_paths: list[Path],
    pptx_path: Path | None,
) -> Path:
    preview_path = result_dir / PPTX_PREVIEW_FILENAME
    payload = {
        "pptx_path": str(pptx_path) if pptx_path else None,
        "slides": _ppt_slide_payload(explanation, analysis_result, chart_paths),
    }
    preview_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return preview_path


def _ppt_slide_payload(
    explanation: dict[str, Any],
    analysis_result: dict[str, Any],
    chart_paths: list[Path],
) -> list[dict[str, Any]]:
    outline = _as_outline(explanation.get("ppt_outline")) if isinstance(explanation, dict) else []
    slides: list[dict[str, Any]] = []
    if outline:
        for index, item in enumerate(outline[:8], start=1):
            chart = str(item.get("chart") or "")
            if not chart and index <= len(chart_paths):
                chart = str(chart_paths[index - 1])
            slides.append(
                {
                    "page": index,
                    "title": str(item.get("title") or f"第 {index} 页"),
                    "subtitle": str(item.get("subtitle") or item.get("section_label") or "分析结论"),
                    "bullets": _as_string_list(item.get("bullets")),
                    "chart": chart,
                }
            )
        if slides:
            return slides

    summary = str(explanation.get("summary") or analysis_result.get("summary") or "数据分析结果") if isinstance(explanation, dict) else "数据分析结果"
    findings = _as_string_list(explanation.get("key_findings")) if isinstance(explanation, dict) else _build_generic_findings(analysis_result)
    recommendations = _as_string_list(explanation.get("recommendations")) if isinstance(explanation, dict) else []
    limitations = _as_string_list(explanation.get("limitations")) if isinstance(explanation, dict) else []
    slides = [
        {
            "page": 1,
            "title": "分析摘要",
            "subtitle": "核心结论总览",
            "bullets": [summary],
            "chart": str(chart_paths[0]) if chart_paths else "",
        },
        {
            "page": 2,
            "title": "关键发现",
            "subtitle": "数据中值得优先关注的信号",
            "bullets": findings[:6] or ["暂无关键发现。"],
            "chart": str(chart_paths[0]) if chart_paths else "",
        },
        {
            "page": 3,
            "title": "建议动作",
            "subtitle": "面向业务跟进的下一步",
            "bullets": recommendations[:6] or ["建议结合业务背景进一步核对。"],
            "chart": str(chart_paths[1]) if len(chart_paths) > 1 else "",
        },
    ]
    if limitations:
        slides.append(
            {
                "page": 4,
                "title": "限制说明",
                "subtitle": "结论适用范围与数据边界",
                "bullets": limitations[:6],
                "chart": str(chart_paths[2]) if len(chart_paths) > 2 else "",
            }
        )
    return slides


def _slide_chart_path(slide_data: dict[str, Any], chart_paths: list[Path]) -> Path | None:
    raw = str(slide_data.get("chart") or "").strip()
    if raw:
        path = _resolve_chart_path(raw, chart_paths[0].parent if chart_paths else Path.cwd())
        if path is not None:
            return path
        raw_name = Path(raw.replace("\\", "/")).name
        for chart_path in chart_paths:
            if chart_path.name == raw_name:
                return chart_path
    return chart_paths[0] if chart_paths else None

